#===============================================================================

from pathlib import Path
from tqdm import tqdm
import json
import pandas as pd
import ast
import logging as log
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

from operator import itemgetter
import torch
from sentence_transformers import SentenceTransformer, util
from SPARQLWrapper import SPARQLWrapper2
from xml.dom import minidom
import re
import itertools
import rdflib

#===============================================================================

from mapmaker import MapMaker
from mapknowledge import KnowledgeStore

#===============================================================================

class PathError(Exception):
    pass

#===============================================================================

# BIOBERT = 'gsarti/biobert-nli'
BIOBERT = 'dmis-lab/biobert-v1.1'
biobert_model = SentenceTransformer(BIOBERT)

#===============================================================================

### General terms identification

tqdm.pandas()
namespaces = {
    'UBERON': 'http://purl.obolibrary.org/obo/UBERON_',
    'ILX': 'http://uri.interlex.org/base/ilx_',
    'FMA': 'http://purl.org/sig/ont/fma/fma'
}

def get_curie(url: str):
    for namespace, preff in namespaces.items():
        if url.startswith(preff):
            return url.replace(preff, namespace + ':')
    return url


#===============================================================================

### Missing nodes and rendered identification

class FlatMapCheck:
    def __init__(self, args):

        # intialisation
        self.__manifest_file = Path(args.manifest_file)
        with open(self.__manifest_file, 'r') as f:
            self.__manifest = json.load(f)
        self.__species = self.__manifest.get('id')
        self.__output_dir = Path(args.output_dir)/f"{self.__species}_output"
        self.__output_dir.mkdir(parents=True, exist_ok=True)
        self.__artefact_dir = Path(args.artefact_dir)
        self.__clean_connectivity = args.clean_connectivity
        self.__align_general = args.align_general
        self.__k = args.k
        if args.sckan_version is not None:
            self.__sckan_version = args.sckan_version
        else:
            self.__sckan_version = self.__manifest.get('sckan-version')
        self.__store = KnowledgeStore(sckan_version=self.__sckan_version, use_npo=True, sckan_provenance=True,store_directory=self.__artefact_dir)

        # setup npo_graph to check term ancestor
        NPO_TURTLE = f'https://raw.githubusercontent.com/SciCrunch/NIF-Ontology/{self.__sckan_version}/ttl/npo.ttl'
        self.__npo_graph = rdflib.Graph()
        self.__npo_graph.parse(NPO_TURTLE, format='turtle')


        # delete artefact files when clean_connectivity
        if self.__clean_connectivity:
            for item in self.__artefact_dir.iterdir():
                if item.is_file():
                    item.unlink()

        # loading all npo connectivities
        print('Load NPO connectivities ...')
        self.__load_npo_connectivities()

        # loading already identified nodes
        self.__map_node_name_file = self.__artefact_dir/'map_node_name.json'
        self.__map_node_name = {}
        if self.__map_node_name_file.exists():
            with open(self.__map_node_name_file, 'r') as f:
                self.__map_node_name = ast.literal_eval(json.load(f))

        # load knowledgebase
        self.__knowledge_file = self.__artefact_dir/'knowledgebase.json'
        self.__knowledge = {}
        if self.__knowledge_file.exists():
            with open(self.__knowledge_file, 'r') as f:
                self.__knowledge = json.load(f)

        # load existing connectivity_terms
        self.__load_connectivity_terms()

        # get map_log of a particular AC/FC
        print('Generate flatmap ...')
        self.__map_log = self.__generate_flatmap()

        # load terms and embeddings of a particular AC/FC
        print('Load terms and embeddings of flatmap ...')
        self.__terms, self.__term_embeddings = self.__load_flatmap_terms()

        # load already identified ancestors
        print('Load ancestors ...')
        self.__map_ancestor = self.__load_ancestors()

        # load connectivity_terms
        print('Load connectivity_terms ...')

    def __load_connectivity_terms(self):
        self.__connectivity_terms = {}
        self.__flatmap_aliases = {}
        if (connectivity_terms_file:=self.__manifest_file.parent/self.__manifest.get('connectivityTerms','None')).exists():
            with open(connectivity_terms_file, 'r') as f:
                conn_terms = json.load(f)
                for ct in conn_terms:
                    ct_id = ct['id']
                    if isinstance(ct['id'], list):
                        ct_id = tuple([ct_id[0], tuple(ct_id[1])])
                    ct_aliases = list(set([tuple([alias[0], tuple(alias[1])]) if isinstance(alias, list) else alias for alias in ct['aliases']]))
                    if ct_id not in self.__connectivity_terms:
                        self.__connectivity_terms[ct_id] = {
                            'id': ct_id,
                            'name': ct.get('name', ''),
                            'aliases': ct_aliases
                        }
                    else:
                        self.__connectivity_terms[ct_id]['aliases'] = list(set(self.__connectivity_terms[ct_id]['aliases']+ct_aliases))
                    for ct_alias in ct_aliases:
                        if ct_alias in self.__flatmap_aliases:
                            if self.__flatmap_aliases[ct_alias] != ct_id:
                                log.error(f"{ct_alias} id mapped multiple nodes to [{self.__flatmap_aliases[ct_alias]}, {ct_id}]")
                        self.__flatmap_aliases[ct_alias] = ct_id

    def __load_npo_connectivities(self):
        # load all connectivities from npo
        # currently we use NPO rather than scicrunch due to neuron-nlp availability
        phenotypes = {
            ('ilxtr:SensoryPhenotype',): 'sensory',
            ('ilxtr:IntrinsicPhenotype',): 'intracardiac',
            ('ilxtr:ParasympatheticPhenotype', 'ilxtr:PostGanglionicPhenotype'): 'post-ganglionic parasympathetic',
            ('ilxtr:PostGanglionicPhenotype', 'ilxtr:SympatheticPhenotype'): 'post-ganglionic sympathetic',
            ('ilxtr:ParasympatheticPhenotype', 'ilxtr:PreGanglionicPhenotype'): 'pre-ganglionic parasympathetic',
            ('ilxtr:PreGanglionicPhenotype', 'ilxtr:SympatheticPhenotype'): 'pre-ganglionic sympathetic',
        }
        biological_sex = self.__manifest.get('biological-sex','')

        npo_connectivities = {}
        for conn in self.__store.connectivity_paths():
            conn_data = self.__store.entity_knowledge(conn)
            if biological_sex=='' or biological_sex==conn_data.get('biologicalSex', '') or conn_data.get('biologicalSex', '')=='':
                npo_connectivities[conn] = conn_data
                if len(conn_phenotype := tuple(sorted(npo_connectivities[conn].get('phenotypes', [])))) > 0:
                    npo_connectivities[conn]['phenotypes'] = phenotypes.get(conn_phenotype, str(conn_phenotype))
        self.__npo_connectivities = npo_connectivities

    def __generate_flatmap(self):
        log_file = self.__artefact_dir/(f"{self.__species}.log")
        if log_file.exists():
            log_file.unlink()
        options = {
            'source': self.__manifest_file.as_posix(),
            'output': self.__artefact_dir.as_posix(),
            'ignoreGit': True,
            'debug': True,
            'logFile': log_file.as_posix(),
            'cleanConnectivity': self.__clean_connectivity,
            'force': True,
            'sckanVersion': self.__sckan_version
        }
        mapmaker = MapMaker(options)
        mapmaker.make()
        map_log = self.__load_log_file(log_file=log_file)
        return map_log

    def __load_flatmap_terms(self):
        ### Loading anatomical map
        #anatomical_terms =
        # {'urinary_1': {'term': 'UBERON:0001008', 'name': 'renal system'},
        #  'urinary_2': {'term': 'UBERON:0001255', 'name': 'urinary bladder'},
        #  ...
        # }

        terms = {}
        if self.__manifest.get('kind', '') != 'functional':
            ### Loading antomicalMap
            anatomical_file = self.__manifest_file.parent/self.__manifest.get('anatomicalMap')
            with open(anatomical_file, 'r') as f:
                anatomical_terms = json.load(f)

            ### Loading property and stored in anatomical term
            # load property
            property_file = self.__manifest_file.parent/self.__manifest.get('properties')
            with open(property_file, 'r') as f:
                properties = json.load(f)
            # load from features key
            for key, val in properties['features'].items():
                if (_model:=val.get('models')) is not None:
                    anatomical_terms[key] = {'term': _model, 'name': val.get('name')}
                elif (_class:=val.get('class')) is not None:
                    if (_anat:=anatomical_terms.get(_class)) is not None:
                        if 'term' not in _anat and 'models' in _anat:
                            _anat['term'] = _anat['models']
                        anatomical_terms[key] = _anat
                    else:
                        pass
                        # print(key) # probably unused anatomy  -- no model or term

            # load from features networks->centrelines for those having models
            for network in properties['networks'][0]['centrelines']:
                if (_model:=network.get('models')) is not None:
                    anatomical_terms[network['id']] = {'term':network.get('models')}

            ## Complete anatomical_terms with no name and check it's concistency
            anaterms = {}
            for term_id in tqdm(set([anaterm['term'] for anaterm in anatomical_terms.values()])):
                if term_id not in anaterms:
                    anaterms[term_id] = self.__get_term_label(term_id)

            ## Select anaterms that available in svg only
            ## Get all id used in csv file
            svg_file = self.__manifest_file.parent/(self.__manifest.get('sources')[0].get('href'))
            doc = minidom.parse(str(svg_file))  # parseString also exists
            svg_used_ids = [path.firstChild.nodeValue[path.firstChild.nodeValue.index('id(')+3:path.firstChild.nodeValue.index(')')].strip() for path in doc.getElementsByTagName('title') if 'id(' in path.firstChild.nodeValue]
            doc.unlink()

            ### Filter anaterms that only available in svg
            # Get terms_ids and term_names
            for idx in set(svg_used_ids) & set(anatomical_terms.keys()):
                term_id = anatomical_terms[idx]['term']
                if term_id not in terms:
                    terms[term_id] = anaterms.get(term_id, term_id).lower()
        else: # handling functional connectivity
            ### Load term in pptx:
            def get_pptx_group_text(shapes):
                text = []
                for shape in shapes:
                    if shape.has_text_frame:
                        if len(term_id:=shape.text.strip().lower()) > 0:
                            text += [term_id]
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        text += get_pptx_group_text(shape.shapes)
                return text

            pptx_text = []
            for file_dict in self.__manifest.get('sources', []):
                pptx_file = self.__manifest_file.parent/file_dict.get('href')
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    pptx_text += get_pptx_group_text(slide.shapes)

            ### Load annotation
            annotation_file = self.__manifest_file.parent/self.__manifest.get('annotation','')
            with open(annotation_file, 'r') as f:
                annotations = json.load(f)
            for _, anatomy_list in annotations.items():
                # if term_type == 'Systems': continue
                for anatomy in anatomy_list:
                    term_id = anatomy.get('Model', anatomy.get('Models', ''))
                    annotation_name = anatomy.get('FTU Name', anatomy.get('Nerve Name', anatomy.get('Organ Name', anatomy.get('System Name', anatomy.get('Vessel Name','')))))
                    if len(term_id.strip()) > 0 and annotation_name.lower() in pptx_text:
                        if (label:=anatomy.get('Label')) is None:
                            label = self.__get_term_label(term_id)
                        terms[term_id] = label.lower()

        ## generate term embedding
        term_embeddings = biobert_model.encode(list(terms.values()), convert_to_tensor=True)
        return terms, term_embeddings

    def __get_term_label(self, term_id):
        if term_id in self.__knowledge:
            return self.__knowledge[term_id]
        self.__knowledge[term_id] = self.__store.label(term_id) # prioritise to npo
        return self.__knowledge[term_id]

    def __get_node_name(self, node):
        name = []
        if node[0] is not None:
            name += [self.__get_term_label(node[0])]

        for n in node[1]:
            if n is not None:
                loc = self.__get_term_label(n)
                name += [loc]
        self.__map_node_name[node] = (name[0], tuple(name[1:] if len(name)>1 else ()))
        return name

    def __get_ancestor(self, term):
        if 'http' in term: return []
        prefix = """
            PREFIX ILX: <http://uri.interlex.org/base/ilx_>
            PREFIX UBERON: <http://purl.obolibrary.org/obo/UBERON_>
            PREFIX TEMP: <http://uri.interlex.org/temp/uris/>
            """
        sparql = f"""
            {prefix}
            SELECT ?object
                WHERE {{
                {term} rdfs:subClassOf ?restriction .
                ?restriction a owl:Restriction ;
                                owl:someValuesFrom ?object .
                }}
            """
        if len(results:=self.__npo_graph.query(sparql)) > 0:
            return [get_curie(str(row.object)) for row in results]
        sparql = f"""
            {prefix}
            SELECT ?object
                WHERE {{
                {term} rdfs:subClassOf ?object .
                }}
            """
        return [get_curie(str(row.object)) for row in self.__npo_graph.query(sparql)]

    def __load_ancestors(self):
        map_ancestor_file = self.__artefact_dir/'map_ancestor.json'
        map_ancestor = {}
        if map_ancestor_file.exists():
            with open(map_ancestor_file, 'r') as f:
                map_ancestor = json.load(f)
        for conn in self.__npo_connectivities.values():
            for edge in conn['connectivity']:
                for term in [edge[0][0]] + list(edge[0][1]) + [edge[1][0]] + list(edge[1][1]):
                    if term not in map_ancestor:
                        level_1_results = {p:1 for p in self.__get_ancestor(term)}
                        level_2_results = {}
                        for p_term in level_1_results.keys():
                            level_2_results = {**level_2_results, **{p:2 for p in self.__get_ancestor(p_term)}}
                        map_ancestor[term] = {**level_1_results, **level_2_results}
        with open(map_ancestor_file, 'w') as f:
            json.dump(map_ancestor, f)
        return map_ancestor

    #===========================================================================

    def __select_ancestor(self, term, cutout=0.1, num_return=1):
        label = self.__get_term_label(term).lower()
        label_emb = biobert_model.encode(label, convert_to_tensor=True)

        ancestors = self.__map_ancestor.get(term, {})
        ancestors = {k:ancestors[k] for k in set(self.__terms.keys()) & set(ancestors.keys())}
        if len(ancestors) > 0:
            target_terms = list(ancestors.keys())
            target_embs = [ self.__term_embeddings[list(self.__terms.keys()).index(k)] for k in target_terms]
            target_embs = torch.stack(target_embs)
            cos_scores = util.cos_sim(label_emb, target_embs)[0]
            top_results = torch.topk(cos_scores, k=len(target_embs))
            results = []
            for score, idx in zip(top_results[0], top_results[1]):
                if score < cutout:
                    break
                term_id = target_terms[idx]
                results += [(term_id, self.__terms[term_id], score.item()+ancestors[term_id])]
            results = sorted(results, key=itemgetter(2))
            return results[:num_return]

        return []

    def __parents(self, x):
        parents = []

        # check the main term
        if x[0] in self.__terms:
            parents = [x[0], []]
        else:
            parent = self.__select_ancestor(x[0])
            if len(parent) > 0:
                parents = [parent[0][0], []]
        # check feature terms
        if len(x[1]) > 0:
            ps = []
            for term in x[1]:
                parent = self.__select_ancestor(term)
                ps += [p[0] for p in parent]
            
            if len(ps) > 0:
                if len(parents) == 0:
                    parents = [ps[0], ps[1:]]
                else:
                    if parents[0] in ps:
                        ps.remove(parents[0])
                    parents[1] = ps
        
        return parents if len(parents) > 0 else None

    #===========================================================================

    ### Candidate term alignment

    def __merge_connectivity_terms(self, df_missing):

        # update df_missing, upadate with parents stated in connectivity_terms
        for parent_id, term in self.__connectivity_terms.items():
            for term_id in term['aliases'] if isinstance(term['aliases'], list) else [term['aliases']]:
                m = df_missing.Node == term_id
                if df_missing[m].parents.isna().any():
                    df_missing.loc[m, 'parents'] = pd.Series([parent_id]*m.sum(), index=m[m].index)

        # save new connectivity_terms
        current_alias = {}
        for idx in df_missing[df_missing.parents.notna()].index:
            if (parent:=df_missing.loc[idx].parents) not in current_alias:
                current_alias[parent] = []
            current_alias[parent] += [df_missing.loc[idx].Node]

        new_alias_dict = {
            node_name:{
                'id': node_name,
                'name': '/'.join(self.__get_node_name(node_name)),
                'aliases': aliases
            }
            for node_name, aliases in current_alias.items()
        }

        # align with current connentivity_terms
        for node, item in new_alias_dict.items():
            if node not in self.__connectivity_terms:
                self.__connectivity_terms[node] = item
            else:
                self.__connectivity_terms[node]['aliases'] = list(set(self.__connectivity_terms[node]['aliases'] + item['aliases']))
                self.__connectivity_terms[node]['name'] = item['name']

        with open(self.__output_dir/'connectivity_terms.json', 'w') as f:
            json.dump(list(self.__connectivity_terms.values()), f, indent=4)

        return df_missing

    def __search_term(self, query, k=5):
        query = query.lower()
        query_emb = biobert_model.encode(query, convert_to_tensor=True)
        cos_scores = util.cos_sim(query_emb, self.__term_embeddings)[0]
        top_results = torch.topk(cos_scores, k=k)
        results = []
        for score, idx in zip(top_results[0], top_results[1]):
            results += [(list(self.__terms.keys())[idx], list(self.__terms.values())[idx], score.item())]
        return results

    def __get_candidates(self, name, k):
        candidates = [[st] for st in self.__search_term(name, k)]

        phrase_candidates = []
        if len(phrases:=name.split(' IN ')) > 1:
            term_candidates = [self.__search_term(phrase, k) for phrase in phrases]
            phrase_candidates = list(itertools.product(*term_candidates))

        of_candidates = []
        if len(phrases:= re.split(r' IN | of ',name)) > 1:
            term_candidates = [self.__search_term(phrase, k) for phrase in phrases]
            of_candidates = list(itertools.product(*term_candidates))

        nodes = []
        for candidate in (candidates+phrase_candidates+of_candidates):
            node = list(zip(*candidate))
            node[0] = [key for key, _ in itertools.groupby(node[0])]
            node[0] = (node[0][0], tuple(node[0][1:] if len(node[0])>1 else []))
            node[1] = [key for key, _ in itertools.groupby(node[1])]
            node[2] = sum(node[2])/len(node[2])
            nodes += [node]

        # return sorted list
        sorted_nodes = sorted(nodes, key=lambda x: x[2], reverse=True)
        selected_nodes, tmp = [], set()
        for node in sorted_nodes:
            if node[0] not in tmp:
                selected_nodes += [node]
                tmp.add(node[0])
            if len(selected_nodes) == k: break
        return selected_nodes

    def __align_missing_nodes(self, df_missing, k):
        if df_missing.shape[0] > 0:
            ### load missing NPO nodes in flatmap
            df_missing = pd.DataFrame(df_missing)
            df_missing['Align candidates'] = df_missing['Node Name'].apply(lambda x: self.__get_candidates(x, k))
            df_missing = df_missing.explode('Align candidates')
            df_missing[['Align candidates', 'Candidate name', 'Score']] = df_missing['Align candidates'].apply(pd.Series)
            df_missing['Selected'] = ''
            df_missing['Comments'] = ''
            return df_missing
        return pd.DataFrame(columns=[['Align candidates', 'Candidate name', 'Score', 'Selected', 'Comments']])

    #===========================================================================

    def __get_missing_nodes(self):
        nodes_to_neuron_types = {}
        for k, v in self.__map_log.items():
            for node in v.get('missing_nodes', []):
                nodes_to_neuron_types[node] = nodes_to_neuron_types.get(node, []) + [k]

        print('Organising missing nodes')
        df = pd.DataFrame(columns=['Node', 'Node Name', 'Appear in', 'Notes'])
        for node, k_types in tqdm(nodes_to_neuron_types.items()):
            name = self.__get_node_name(node)
            name = ' IN '.join(name)
            notes = ''
            if node in self.__connectivity_terms:
                aliases = ', '.join([f'{str(alias)}-{self.__get_node_name(alias)}' for alias in self.__connectivity_terms[node]['aliases']])
                notes = f'The node is available in `connextivity_terms.json` as the aliases of {aliases}. Possible problems:\
                    \n - missing a pink dot in the svg file\
                    \n - incorrect alias mapping'
            elif node in self.__flatmap_aliases:
                notes = f'The node is mapped to `{self.__flatmap_aliases[node]}-{self.__connectivity_terms[self.__flatmap_aliases[node]].get('name','')}` in `connextivity_terms.json`. Possible problem:\
                    \n - other nodes on the same path are generalised to the same node.'
            df.loc[len(df)] = [
                node,
                name,
                '\n'.join(list(set(k_types))),
                notes
            ]
        df = df.sort_values('Appear in')
        return df

    def __load_log_file(self, log_file):
        # a function to load log file
        missing_nodes = {} # node:label
        missing_segments = {} # neuron_path:segment
        missing_paths = [] # a mlist of missing path
        map_log = {}
        with open(log_file, 'r') as f:
            while line := f.readline():
                tag_feature = 'Cannot find feature for connectivity node '
                tag_segment = 'Cannot find any sub-segments of centreline for '
                tag_path = 'Path is not rendered at all.'
                if tag_feature in line:
                    feature = line.split(tag_feature)[-1].split(') (')
                    missing_nodes[ast.literal_eval(f'{feature[0]})')] = f'({feature[1]}'.strip()
                elif tag_segment in line:
                    path_id = line[33:].split(': ')[0]
                    if path_id not in missing_segments: missing_segments[path_id] = []
                    missing_segments[path_id] += [line.split(tag_segment)[-1][1:-2]]
                elif tag_path in line:
                    missing_paths += [(path_id:=line[33:].split(': ')[0])]
        for path_id, connectivities in self.__npo_connectivities.items():
            map_log[path_id] = {}
            nodes, edges = set(), set()
            m_edges = set()
            for edge in connectivities['connectivity']:
                node_0 = (edge[0][0], tuple(edge[0][1]))
                node_1 = (edge[1][0], tuple(edge[1][1]))
                flat_nodes = set([edge[0][0]] + list(edge[0][1]) + [edge[1][0]] + list(edge[1][1]))
                nodes.add(node_0)
                nodes.add(node_1)
                edges.add((node_0, node_1))

                # check if nodes are in the missing list
                if node_0 in missing_nodes or node_1 in missing_nodes:
                    m_edges.add((node_0, node_1))
                # check if edge contain missing segment
                if len(set(missing_segments.get(path_id, [])) & flat_nodes) > 1:
                    m_edges.add((node_0, node_1))
            mapped_nodes = set([n if n not in self.__flatmap_aliases else self.__flatmap_aliases[n] for n in nodes])
            m_nodes = mapped_nodes & set(missing_nodes.keys())
            r_nodes = mapped_nodes - set(missing_nodes.keys())
            r_edges = edges - m_edges
            complete = 'Complete' if len(m_nodes)==0 and len(m_edges)==0 else 'Partial'
            notes = ''
            if path_id in missing_paths:
                m_nodes = [n for n in nodes if n in self.__flatmap_aliases]
                r_nodes = []
                m_edges = edges
                r_edges = []
                complete = 'None'
                notes = f'All nodes are generalised to the same term, so this path is not rendered.\nNodes: {nodes}'
            map_log[path_id] = {
                'original_nodes': nodes,
                'missing_nodes': m_nodes,
                'rendered_nodes': r_nodes,
                'original_edges': edges,
                'missing_edges': m_edges,
                'rendered_edges': r_edges,
                'completeness': complete,
                'missing_segment': missing_segments.get(path_id, []),
                'notes': notes
            }

        return map_log

    def __organised_map_log(self):
        # a function to organised data into dataframe and then save it as csv file
        ### complete neuron:
        df = pd.DataFrame(columns=['Neuron NPO', 'Completeness', 'Missing Nodes', 'Missing Node Name', 'Missing Edges', 'Missing Edge Name', 'Missing Segments', 'Missing Segment Name', 'Rendered Edges', 'Rendered Edge Name', 'Notes'])
        keys = [
            'missing_nodes',
            'missing_edges',
            'missing_segments',
            'rendered_edges'
        ]
        print('Organising map log')

        for neuron, value in tqdm(self.__map_log.items()):
            info = {}
            for key in keys:
                info[key] = '\n'.join([str(mn) for mn in list(value.get(key,[]))])

                if len(value.get(key,[]))>0:
                    if key not in ['missing_edges', 'rendered_edges']:
                        names = [self.__map_node_name[node] for node in value.get(key,[])]
                    else:
                        names = []
                        for edge in value.get(key,[]):
                            if edge[0] not in self.__map_node_name: self.__get_node_name(edge[0])
                            if edge[1] not in self.__map_node_name: self.__get_node_name(edge[1])
                            names += [(self.__map_node_name[edge[0]], self.__map_node_name[edge[1]])]
                else:
                    names = ''
                info[key+'_name'] = '\n'.join([str(mnn) for mnn in names])
            info['Notes'] = value['notes']
            df.loc[len(df)] = [neuron, value['completeness']] + list(info.values())

        df = df.sort_values('Completeness', ascending=False)
        return df

    def check_npo_in_flatmap(self):
        # identified the miissing nodes
        df_missing = self.__get_missing_nodes()
        # identify the missing node general terms when is align_general
        df_missing['parents'] = None
        if self.__align_general:
            df_missing['parents'] = df_missing['Node'].progress_apply(lambda x: self.__parents(x))

        # update and merge identified parent with the existing connectivity_terms
        df_missing = self.__merge_connectivity_terms(df_missing)

        # save missing_nodes to a file
        df_missing = df_missing[df_missing.parents.isna() | df_missing.Node.isin(self.__flatmap_aliases)]
        
        # align missing_nodes and safe to file
        df_align = self.__align_missing_nodes(df_missing, self.__k)

        # save rendered_nodes to a file
        df_rendered = self.__organised_map_log()

        # save to excel
        def set_format(writer, df, sheet_name):
            workbook  = writer.book
            wrap_format = workbook.add_format({'text_wrap': True})
            for column in df:
                col_idx = df.columns.get_loc(column)
                writer.sheets[sheet_name].set_column(col_idx, col_idx, 25, wrap_format)

        excel_file = self.__output_dir/f'npo_completeness-{self.__species}.xlsx'
        with pd.ExcelWriter(excel_file, engine='xlsxwriter') as writer:
            df_missing.to_excel(writer, sheet_name='missing', index=False)
            set_format(writer, df_missing, 'missing')
            df_align.to_excel(writer, sheet_name='align', index=False)
            set_format(writer, df_align, 'align')
            df_rendered.to_excel(writer, sheet_name='rendered', index=False)
            set_format(writer, df_rendered, 'rendered')

        with open(self.__map_node_name_file, 'w') as f:
            json.dump(str(self.__map_node_name), f)

        with open(self.__knowledge_file, 'w') as f:
            json.dump(self.__knowledge, f)

#===============================================================================

def main():
    import argparse
    import sys

    parser = argparse.ArgumentParser(description="Checking nodes and edges completeness in the generated flatmap")
    parser.add_argument('--manifest', dest='manifest_file', metavar='MANIFEST', help='Path of flatmap manifest')
    parser.add_argument('--sckan-version', dest='sckan_version', metavar='SCKAN_VERSION', help='SCKAN version to check, e.g. sckan-2024-03-26')
    parser.add_argument('--artefact-dir', dest='artefact_dir', metavar='ARTEFACT_DIR', help='Directory to store artefact files, e.g. generated maps and log file, to check NPO completeness')
    parser.add_argument('--output-dir', dest='output_dir', metavar='OUTPUT_DIR', help='Directory to store the check results')
    parser.add_argument('--clean-connectivity', dest='clean_connectivity', action='store_true', help='Run mapmaker as a clean connectivity (optional)')
    parser.add_argument('--align-general-term', dest='align_general', action='store_true', help='Find general terms of the missing nodes to align. This is useful for FC alignment')
    parser.add_argument('--k', dest='k', help='The number of generated candidates for earch missing nodes', default=5)

    try:
        args = parser.parse_args()
        flatmap_ckeck = FlatMapCheck(args)
        flatmap_ckeck.check_npo_in_flatmap()
    except PathError as error:
        sys.stderr.write(f'{error}\n')
        sys.exit(1)
    sys.exit(0)

#===============================================================================

if __name__ == '__main__':
    main()

#===============================================================================

# Documentation
# after running the environment with `poetry shell`, run this script

# Running:
# python ./npo_check.py --manifest `manifest file` \
#                       --artefact-dir `any directory to store generated files` \
#                       --output-dir 'a directory to save csv file'

# Options:
# --manifest `manifest file`
# --artefact-dir `any directory to store generated files`
# --output-dir `a directory to save csv file`
# --clean-connectivity
# --align-general-term
# --k `integer value > 0`
# --sckan-version `the version of SCKAN, e.g. sckan-2024-03-26`


# Results are stored in --output-dir/--species/ directory
#   - npo_{species}_completeness.xlsx, There are 3 sheets inside:
#       - missing
#        - align
#       - rendered
#   - connectivity_terms.json,
#       - combines the existing connectivity_terms.json with the general terms found
#       - later this file can be used as a new connectivity_terms.json in AC and FC
#       - final connectivity_terms should be combined with manually align terms using npo_alias.py

# in order to generate candidate alignment, run npo_align.py file
