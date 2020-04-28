#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from copy import deepcopy
import pptx

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

LAYOUT_BLANK_SLIDE = 6

#===============================================================================

XPATH_GROUP_SPPR        = './p:grpSpPr'

XPATH_SHAPE_BY_ID_GROUP = './p:sp/p:nvSpPr/p:cNvPr[@id={}]'
XPATH_SHAPE_BY_ID_SHAPE = './p:nvSpPr/p:cNvPr[@id={}]'

XPATH_CXN_BY_ID_GROUP   = './p:cxnSp/p:nvCxnSpPr/p:cNvPr[@id={}]'
XPATH_CXN_BY_ID_SHAPE   = './p:nvCxnSpPr/p:cNvPr[@id={}]'

XPATH_CXN_TYPE          = './p:spPr/a:prstGeom[@prst]'

XPATH_SLIDE_cSld        = './p:cSld'
XPATH_SLIDE_extLst      = './p:cSld/p:extLst'

XPATH_PRS_extLst        = './p:extLst'

#===============================================================================

def valid_name(name):
    if name.startswith('.'):
        for directive in name[1:].split():
            if directive in ['group', 'invisible', 'region']:
                return False
    return True

#===============================================================================

def valid_notes(notes):
    return notes and notes[0] in ['.', '#']

#===============================================================================

def clean_markup(text):
    return ('.' + text[1:]) if (text and text[0] == '#') else text

#===============================================================================

def connector_type(name):
    return (MSO_CONNECTOR.STRAIGHT if name == 'line'
       else MSO_CONNECTOR.ELBOW    if name == 'bentConnector3'
       else MSO_CONNECTOR.CURVE    if name == 'curvedConnector3'
       else MSO_CONNECTOR.MIXED)

#===============================================================================

class CleanPresentation(object):
    def __init__(self, width, height):
        self._prs = pptx.Presentation()
        self._prs.slide_width = width
        self._prs.slide_height = height
        self._blank_layout = self._prs.slide_layouts[LAYOUT_BLANK_SLIDE]
        self._current_group = None
        self._group_stack = []
        self._clean_slide = None

    def save(self, output_file):
        if len(self._group_stack):
            raise ValueError('Unclosed group...')

        # Add `p:extLst` to new presentation
        self._prs.element.append(self._source.element.xpath(XPATH_PRS_extLst)[0])

        self._prs.save(output_file)

    def add_slide(self, slide):
        self._clean_slide = self._prs.slides.add_slide(self._blank_layout)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if valid_notes(notes):
                clean_notes_slide = self._clean_slide.notes_slide
                clean_notes_slide.notes_text_frame.text = clean_markup(notes)

        self.start_shapes_()
        self.add_shapes_(slide.shapes)
        self.end_shapes_(slide)

        xml = open('clean_slide.xml', 'w')
        xml.write(self._clean_slide.element.xml)
        xml.close()
        # Add `p:extLst` to new `cSld` element  ### ????
        new_csld_element = self._clean_slide.element.xpath(XPATH_SLIDE_cSld)[0]
        new_csld_element.append(slide.element.xpath(XPATH_SLIDE_extLst)[0])

    def start_shapes_(self):
        self._current_group = self._clean_slide
        self._group_stack = []

    def end_shapes_(self, slide):
        # replace <p:grpSpPr> element of the slide's shapes group
        clean_shapes_element = self._clean_slide.shapes.element
        clean_shapes_element.replace(clean_shapes_element.xpath(XPATH_GROUP_SPPR)[0],
                                     slide.shapes.element.xpath(XPATH_GROUP_SPPR)[0])

    def add_shapes_(self, shapes):
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or shape.shape_type == MSO_SHAPE_TYPE.PICTURE
             or isinstance(shape, pptx.shapes.connector.Connector)):
                if valid_name(shape.name):
                # parse and filter out region etc
                    self.append_shape_(shape)

            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.start_group_()
                self.add_shapes_(shape.shapes)
                self.end_group_(shape)

            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                pass

            else:
                print('Unknown shape type {}, "{}"...'.format(str(shape.shape_type), shape.name))

    def append_shape_(self, shape):
        # We add a new shape and then replace
        # its xml element with the shape being appended
        if isinstance(shape, pptx.shapes.connector.Connector):
            cxn_name = shape.element.xpath(XPATH_CXN_TYPE)[0].get('prst')
            new_shape = self._current_group.shapes.add_connector(connector_type(cxn_name),
                                                                 shape.begin_x, shape.begin_y,
                                                                 shape.end_x, shape.end_y)
            new_id_element_xpath = XPATH_CXN_BY_ID_GROUP.format(new_shape.shape_id)
            old_id_element_xpath = XPATH_CXN_BY_ID_SHAPE.format(shape.shape_id)
        else:
            new_shape = self._current_group.shapes.add_shape(shape.shape_type,
                                                             shape.left, shape.top,
                                                             shape.width, shape.height)
            new_id_element_xpath = XPATH_SHAPE_BY_ID_GROUP.format(new_shape.shape_id)
            old_id_element_xpath = XPATH_SHAPE_BY_ID_SHAPE.format(shape.shape_id)

        new_id_element = self._current_group.shapes.element.xpath(new_id_element_xpath)[0]
        new_element = new_id_element.getparent().getparent()

        shape_element = deepcopy(shape.element)
        shape_id_element = shape_element.xpath(old_id_element_xpath)[0]
        shape_id_element.set('id', str(new_shape.shape_id))
        self._current_group.shapes.element.replace(
            new_element,
            shape_element
            )

    def start_group_(self):
        self._group_stack.append(self._current_group)
        self._current_group = self._current_group.shapes.add_group_shape()

    def end_group_(self, group):
        self._current_group.name = group.name
        # replace <p:grpSpPr> element of group
        clean_shapes_element = self._current_group.shapes.element
        clean_shapes_element.replace(clean_shapes_element.xpath(XPATH_GROUP_SPPR)[0],
                                     group.shapes.element.xpath(XPATH_GROUP_SPPR)[0])
        self._current_group = self._group_stack.pop()

#===============================================================================

class Presentation(object):
    def __init__(self, source_file):
        self._prs = pptx.Presentation(source_file)

    def clean(self, output_file):
        cleaned = CleanPresentation(self._prs.slide_width, self._prs.slide_height)
        for slide in self._prs.slides:
            cleaned.add_slide(slide)
        cleaned.save(output_file)

#===============================================================================

def clean(source, target):
#=========================
    presentation = Presentation(source)
    presentation.clean(target)

#===============================================================================

if __name__ == '__main__':
    clean('map_sources/Rat_flatmap_annotation_test_Dave.pptx',
          'map_sources/cleaned.pptx')

#===============================================================================
